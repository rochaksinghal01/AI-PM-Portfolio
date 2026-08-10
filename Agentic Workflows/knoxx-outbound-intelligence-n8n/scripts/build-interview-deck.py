#!/usr/bin/env python3
"""Build the Knoxx interview PPTX and a PNG architecture asset."""

from pathlib import Path
from textwrap import wrap

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
DOCS = ROOT / "docs"
OUT = DOCS / "Knoxx-Outbound-Intelligence-Interview.pptx"
ARCH_PNG = DOCS / "architecture-clean.png"

NAVY = "111827"
NAVY_2 = "1E293B"
SLATE = "475569"
MUTED = "94A3B8"
PAPER = "F8FAFC"
WHITE = "FFFFFF"
ORANGE = "F97316"
ORANGE_PALE = "FFF7ED"
VIOLET = "7C3AED"
VIOLET_PALE = "F5F3FF"
GREEN = "059669"
GREEN_PALE = "ECFDF5"
BLUE = "2563EB"
BLUE_PALE = "EFF6FF"
AMBER = "D97706"
AMBER_PALE = "FFFBEB"
RED = "DC2626"
RED_PALE = "FEF2F2"


def rgb(value):
    return RGBColor.from_string(value)


def font(size, bold=False):
    candidates = [
        "/System/Library/Fonts/Supplemental/Arial.ttf",
        "/System/Library/Fonts/Supplemental/Arial Bold.ttf" if bold else "",
        "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
    ]
    if bold:
        candidates.insert(0, "/System/Library/Fonts/Supplemental/Arial Bold.ttf")
        candidates.insert(1, "/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf")
    for candidate in candidates:
        if candidate and Path(candidate).exists():
            return ImageFont.truetype(candidate, size)
    return ImageFont.load_default()


def draw_architecture():
    image = Image.new("RGB", (1920, 1080), "#F8FAFC")
    draw = ImageDraw.Draw(image)
    title = font(46, True)
    subtitle = font(23)
    heading = font(16, True)
    box_title = font(21, True)
    body = font(16)
    draw.text((80, 56), "Knoxx Outbound Intelligence", fill="#111827", font=title)
    draw.text((80, 116), "One domain → auditable account decision → safe outreach → account-level stop", fill="#475569", font=subtitle)

    def arrow(x1, y1, x2, y2, color="#64748B", width=5):
        draw.line((x1, y1, x2, y2), fill=color, width=width)
        if abs(x2-x1) > abs(y2-y1):
            points = [(x2, y2), (x2-16 if x2>x1 else x2+16, y2-10), (x2-16 if x2>x1 else x2+16, y2+10)]
        else:
            points = [(x2, y2), (x2-10, y2-16 if y2>y1 else y2+16), (x2+10, y2-16 if y2>y1 else y2+16)]
        draw.polygon(points, fill=color)

    def box(x, y, w, h, title_text, lines, fill, outline):
        draw.rounded_rectangle((x, y, x+w, y+h), radius=20, fill=fill, outline=outline, width=3)
        draw.text((x+24, y+20), title_text, fill="#111827", font=box_title)
        for i, line in enumerate(lines):
            draw.text((x+24, y+58+i*25), line, fill="#475569", font=body)

    draw.text((80, 190), "PRODUCT + CONTROL", fill="#64748B", font=heading)
    box(80, 220, 245, 125, "Sales operator", ["Domain + notes", "Review and approve"], "#FFFFFF", "#C7D2FE")
    box(380, 220, 245, 125, "React / Lovable", ["Pipeline + report", "Publishable key only"], "#FFFFFF", "#A7F3D0")
    box(680, 220, 245, 125, "Supabase Edge", ["Auth + protected APIs", "Returns run_id"], "#ECFDF5", "#34D399")
    box(980, 220, 245, 125, "Postgres + RLS", ["Source of truth", "Runs + evidence + state"], "#FFFFFF", "#94A3B8")
    arrow(325, 282, 372, 282); arrow(625, 282, 672, 282); arrow(925, 282, 972, 282)

    draw.text((80, 420), "N8N WORKFLOWS", fill="#64748B", font=heading)
    stages = [
        ("WF01", "Intake", "#FFF7ED", "#FB923C"),
        ("WF02", "Evidence + fit", "#F5F3FF", "#A78BFA"),
        ("WF04", "Committee", "#EFF6FF", "#60A5FA"),
        ("WF05", "Draft + safe send", "#FDF2F8", "#F472B6"),
        ("WF06", "Events + stop", "#ECFEFF", "#22D3EE"),
    ]
    xs = [80, 390, 700, 1010, 1320]
    for x, (wf, label, fill, outline) in zip(xs, stages):
        box(x, 450, 245, 125, f"{wf} · {label}", ["Persist checkpoint", "Evaluate + explain"], fill, outline)
    for x in [325, 635, 945, 1255]:
        arrow(x, 512, x+57, 512)

    draw.text((80, 650), "BOUNDED PROVIDERS", fill="#64748B", font=heading)
    providers = [("Firecrawl", "web + PDFs"), ("OpenAI", "bounded agents"), ("Apollo", "search + enrich"), ("Gmail", "allowlisted send")]
    for i, (name, line) in enumerate(providers):
        box(80+i*300, 680, 245, 105, name, [line], "#FFFFFF", ["#FDBA74", "#C4B5FD", "#93C5FD", "#86EFAC"][i])
    arrow(325, 732, 372, 732); arrow(625, 732, 672, 732); arrow(925, 732, 972, 732)
    box(1320, 680, 520, 105, "Evaluation harness", ["Frozen fixtures + deterministic metrics", "+ independent Gemini judge"], "#FFFBEB", "#FBBF24")

    draw.rounded_rectangle((80, 890, 1760, 990), radius=22, fill="#111827")
    draw.text((112, 918), "1,000-domain target state", fill="#FFFFFF", font=box_title)
    draw.text((112, 954), "202 + run_id  ·  durable queue  ·  bounded workers  ·  cache  ·  retries/DLQ  ·  observability", fill="#CBD5E1", font=body)
    image.save(ARCH_PNG, quality=95)


def add_text(slide, text, x, y, w, h, size=18, color=NAVY, bold=False, align=PP_ALIGN.LEFT, font_name="Aptos", margin=0.04):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear(); tf.margin_left = Inches(margin); tf.margin_right = Inches(margin); tf.margin_top = Inches(margin); tf.margin_bottom = Inches(margin)
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.text = text; p.alignment = align
    for run in p.runs:
        run.font.name = font_name; run.font.size = Pt(size); run.font.bold = bold; run.font.color.rgb = rgb(color)
    return box


def rect(slide, x, y, w, h, fill=WHITE, line="E2E8F0", radius=True):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid(); shape.fill.fore_color.rgb = rgb(fill)
    shape.line.color.rgb = rgb(line)
    return shape


def base_slide(prs, title, kicker=None, dark=False):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill; bg.solid(); bg.fore_color.rgb = rgb(NAVY if dark else PAPER)
    if kicker:
        add_text(slide, kicker.upper(), .65, .35, 7.5, .28, 10, ORANGE if dark else VIOLET, True)
    add_text(slide, title, .65, .68, 12.0, .72, 27, WHITE if dark else NAVY, True)
    return slide


def footer(slide, number, source=None, dark=False):
    color = "CBD5E1" if dark else MUTED
    add_text(slide, "KNOXX OUTBOUND INTELLIGENCE", .65, 7.18, 4.5, .2, 8, color, True)
    if source:
        add_text(slide, source, 4.4, 7.15, 7.8, .26, 7.5, color, False, PP_ALIGN.RIGHT)
    add_text(slide, f"{number:02d}", 12.45, 7.13, .35, .25, 9, ORANGE, True, PP_ALIGN.RIGHT)


def add_bullets(slide, items, x, y, w, h, size=17, color=NAVY, bullet_color=ORANGE, gap=0.48):
    cursor = y
    for item in items:
        rect(slide, x, cursor+.08, .10, .10, bullet_color, bullet_color)
        add_text(slide, item, x+.22, cursor, w-.22, gap, size, color)
        cursor += gap


def add_card(slide, x, y, w, h, title, body, fill=WHITE, accent=VIOLET, title_size=16, body_size=12.5):
    rect(slide, x, y, w, h, fill, "E2E8F0")
    rect(slide, x, y, .07, h, accent, accent, False)
    add_text(slide, title, x+.23, y+.18, w-.4, .32, title_size, NAVY, True)
    add_text(slide, body, x+.23, y+.58, w-.4, h-.7, body_size, SLATE)


def add_table(slide, headers, rows, x, y, widths, row_h=.56, font_size=11):
    total = sum(widths)
    current = x
    for header, width in zip(headers, widths):
        rect(slide, current, y, width, row_h, NAVY_2, NAVY_2, False)
        add_text(slide, header, current+.08, y+.13, width-.16, .28, font_size, WHITE, True)
        current += width
    for i, row in enumerate(rows):
        current = x; yy = y + row_h*(i+1); fill = WHITE if i % 2 == 0 else "F1F5F9"
        for cell, width in zip(row, widths):
            rect(slide, current, yy, width, row_h, fill, "E2E8F0", False)
            add_text(slide, str(cell), current+.08, yy+.10, width-.16, row_h-.14, font_size, NAVY if i == 0 else SLATE, False)
            current += width
    return total


def build_deck():
    draw_architecture()
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # 1 title
    slide = base_slide(prs, "Knoxx Outbound Intelligence", dark=True)
    add_text(slide, "AI-powered account research, ingredient qualification and safe outreach", .7, 1.62, 10.7, .8, 23, "CBD5E1")
    rect(slide, .7, 2.75, 2.55, .42, ORANGE, ORANGE)
    add_text(slide, "TECHNICAL INTERVIEW CASE STUDY", .84, 2.85, 2.28, .2, 9, WHITE, True)
    add_text(slide, "Rochak Singhal", .7, 5.72, 4.0, .4, 18, WHITE, True)
    add_text(slide, "One company per run · evidence before judgment · humans approve outreach", .7, 6.16, 9.5, .35, 14, "94A3B8")
    footer(slide, 1, dark=True)

    # 2 problem
    slide = base_slide(prs, "A decision system—not an email generator", "Product definition")
    add_card(slide, .65, 1.55, 3.85, 2.0, "Research is fragmented", "Websites, PDFs, spreadsheets and contact tools make sales research slow and hard to audit.", WHITE, ORANGE, 18, 14)
    add_card(slide, 4.75, 1.55, 3.85, 2.0, "Qualification comes first", "The system decides whether an account deserves enrichment and outreach before spending credits or rep time.", WHITE, VIOLET, 18, 14)
    add_card(slide, 8.85, 1.55, 3.85, 2.0, "Control remains explicit", "AI proposes judgments. Code owns identity, formulas, thresholds, approval, sending and stop rules.", WHITE, GREEN, 18, 14)
    rect(slide, .65, 4.05, 12.05, 1.4, NAVY, NAVY)
    add_text(slide, "domain", .95, 4.49, 1.0, .3, 16, WHITE, True); add_text(slide, "→", 1.85, 4.44, .35, .35, 21, ORANGE, True)
    stages = ["evidence", "ingredient fit", "committee", "draft", "approval", "engagement", "stop"]
    xx = 2.2
    for i, value in enumerate(stages):
        add_text(slide, value, xx, 4.49, 1.15, .3, 14, "E2E8F0", True)
        xx += 1.38
        if i < len(stages)-1:
            add_text(slide, "→", xx-0.25, 4.45, .30, .35, 18, ORANGE, True)
    add_text(slide, "One account may have many immutable research runs.", .75, 5.75, 7.2, .35, 15, SLATE)
    footer(slide, 2)

    # 3 product outcomes
    slide = base_slide(prs, "Success is better prioritization—not more automated email", "User · job · outcome")
    add_card(slide, .65, 1.45, 3.7, 1.75, "PRIMARY USER", "Salesperson or business-development manager deciding where to spend research and outreach time.", BLUE_PALE, BLUE, 17, 13.5)
    add_card(slide, 4.58, 1.45, 4.0, 1.75, "JOB TO BE DONE", "Decide if an account is worth pursuing, show why, identify who matters, recommend the next action and know when to stop.", VIOLET_PALE, VIOLET, 17, 13.5)
    add_card(slide, 8.82, 1.45, 3.88, 1.75, "NORTH-STAR CANDIDATE", "Sales-accepted qualified accounts per research hour—not email volume.", GREEN_PALE, GREEN, 17, 13.5)
    rows = [
        ("Speed", "Median research minutes per reviewed account"),
        ("Precision", "% qualified accounts accepted by sales"),
        ("Trust", "Citation coverage + unsupported-claim rate"),
        ("Economics", "Cost per completed and accepted account"),
    ]
    add_table(slide, ["PILOT OUTCOME", "MEASURE"], rows, .65, 3.62, [3.0, 6.2], .58, 11.5)
    add_card(slide, 10.12, 3.62, 2.58, 2.9, "GUARDRAILS", "Zero unintended recipients\n\nZero unsupported golden-case claims\n\n100% stop/suppression tests\n\nHuman approval", RED_PALE, RED, 14, 11.3)
    add_text(slide, "Baseline targets with sales in a controlled pilot; do not invent impact numbers in advance.", .75, 6.35, 9.0, .34, 13, VIOLET, True)
    footer(slide, 3)

    # 4 assumptions/questions
    slide = base_slide(prs, "Assumptions made vs. questions sales must answer", "Discovery")
    rows = [
        ("ICP", "Prepared-food manufacturers", "What makes an account unquestionably valuable?"),
        ("Catalogue", "6 synthetic SKUs", "Which SKUs, capacity, MOQ and claims are approved?"),
        ("Demand", "Meals × share × kg/meal", "What uncertainty and volume floor are acceptable?"),
        ("Committee", "Procurement, NPD, ops", "Who champions, approves, buys and blocks?"),
        ("Outreach", "Days 0 / 3 / 7 / 12", "Which CTA, proof and stop conditions are approved?"),
    ]
    add_table(slide, ["STAGE", "DEMO ASSUMPTION", "PRODUCTION QUESTION"], rows, .65, 1.48, [2.0, 3.6, 6.45], .78, 12)
    add_text(slide, "The product is configurable only after sales policy is explicit and versioned.", .75, 6.25, 10.5, .36, 15, VIOLET, True)
    footer(slide, 4)

    # 5 architecture
    slide = base_slide(prs, "Q1 · Architecture walkthrough", "Raw domain to completed database state")
    slide.shapes.add_picture(str(ARCH_PNG), Inches(.65), Inches(1.35), width=Inches(12.05), height=Inches(5.5))
    add_text(slide, "Architecture source: docs/architecture-clean.svg", .75, 6.74, 5.4, .26, 10, VIOLET, True)
    footer(slide, 5)

    # 6 flow
    slide = base_slide(prs, "Eight steps from input to an auditable row", "Data flow")
    items = [
        ("01", "Accept", "Edge Function returns account_id + run_id."),
        ("02", "Canonicalize", "WF01 validates domain and deduplicates."),
        ("03", "Collect", "WF02 crawls bounded web/PDF evidence."),
        ("04", "Interpret", "Agent extracts dishes, scale, pains and matches."),
        ("05", "Calculate", "Code computes ranges, risk and fit tier."),
        ("06", "Prioritize", "WF04 ranks 3–5 contacts; enrichs top 1–2."),
        ("07", "Approve", "WF05 drafts; human approves; safe inbox wins."),
        ("08", "React", "WF06 records events and stops the organization."),
    ]
    for i, (num, title, body) in enumerate(items):
        col=i%2; row=i//2; x=.65+col*6.05; y=1.4+row*1.33
        rect(slide, x, y, 5.78, 1.05, WHITE, "E2E8F0")
        rect(slide, x+.18, y+.18, .56, .56, ORANGE if i<4 else VIOLET, ORANGE if i<4 else VIOLET)
        add_text(slide, num, x+.25, y+.34, .42, .2, 11, WHITE, True, PP_ALIGN.CENTER)
        add_text(slide, title, x+.9, y+.17, 1.3, .25, 15, NAVY, True)
        add_text(slide, body, x+.9, y+.49, 4.55, .34, 11.5, SLATE)
    footer(slide, 6)

    # 7 agent boundary
    slide = base_slide(prs, "AI interprets; deterministic controls decide", "Agent boundary")
    add_card(slide, .65, 1.5, 5.7, 4.85, "AI is used for ambiguity", "• Menu and catalogue language\n\n• Observed vs inferred vs hypothesis\n\n• Ingredient synonym and product relevance\n\n• Job-title and persona relevance\n\n• Persona-specific message framing\n\n• Bounded reply classification", VIOLET_PALE, VIOLET, 20, 15)
    add_card(slide, 6.65, 1.5, 5.7, 4.85, "Code is used for control", "• Domain identity and deduplication\n\n• Source-key and catalogue-ID validation\n\n• Quantity arithmetic and score thresholds\n\n• Approval, scheduling and recipient override\n\n• Suppression, idempotency and state\n\n• Account-level stop rules", GREEN_PALE, GREEN, 20, 15)
    footer(slide, 7)

    # 8 prompting
    slide = base_slide(prs, "Q2 · Prompt engineering: constrain the decision surface", "Reusable prompt anatomy")
    components = ["ROLE", "OBJECTIVE", "CONTEXT", "EVIDENCE", "PROCEDURE", "TOOLS", "GUARDRAILS", "SCHEMA", "EXAMPLES", "SELF-CHECK"]
    descriptions = ["one bounded job", "business decision", "account + rules", "known source keys", "ordered reasoning", "mandatory lookups", "never invent/act", "typed JSON", "supported + unknown", "citations + nulls"]
    for i, (c,d) in enumerate(zip(components, descriptions)):
        col=i%5; row=i//5; x=.65+col*2.43; y=1.65+row*1.55
        rect(slide,x,y,2.18,1.17,WHITE,"E2E8F0")
        add_text(slide,f"{i+1:02d}",x+.14,y+.13,.38,.2,10,ORANGE,True)
        add_text(slide,c,x+.14,y+.42,1.9,.24,13,NAVY,True)
        add_text(slide,d,x+.14,y+.76,1.9,.2,10.5,SLATE)
    rect(slide,.65,5.1,12.05,1.03,NAVY,NAVY)
    add_text(slide,"Key separation",.92,5.34,1.55,.26,15,ORANGE,True)
    add_text(slide,"The model proposes evidence-backed inputs. The workflow calculates volume, score and tier.",2.45,5.31,9.75,.34,16,WHITE,True)
    footer(slide, 8)

    # 9 model routing
    slide = base_slide(prs, "Model selection is a routing decision", "Quality × latency × cost")
    rows = [
        ("Retrieval", "Firecrawl", "Preserve URLs/passages before interpretation"),
        ("Research", "GPT-5.6 Terra", "Balanced multi-source reasoning + tools"),
        ("Committee", "GPT-5.6 Terra", "Title/organization reasoning"),
        ("Outreach", "GPT-5.6 Luna", "Bounded, high-volume copy"),
        ("Reply", "GPT-5.6 Luna", "Narrow classification taxonomy"),
        ("Judge", "Gemini 3.6 Flash", "Independent provider; fast structured output"),
        ("Escalation", "GPT-5.6 Sol", "Hard cases only; human follows"),
    ]
    add_table(slide,["STAGE","DEFAULT","WHY"],rows,.65,1.45,[2.15,3.15,6.75],.66,11.5)
    add_text(slide,"Promote a replacement only when golden evaluations improve without breaking latency/cost and safety gates.",.75,6.45,11.5,.38,14,VIOLET,True)
    footer(slide, 9, "OpenAI model docs · Gemini model/structured-output docs")

    # 10 eval
    slide = base_slide(prs, "Evaluations are release gates", "Structured parser → objective metrics → independent judge")
    add_card(slide,.65,1.6,3.75,2.05,"1 · Contract","Strict schema catches missing fields, wrong types and malformed JSON.",WHITE,BLUE,18,14)
    add_card(slide,4.78,1.6,3.75,2.05,"2 · Hard metrics","Citation coverage, allowed IDs, tier/persona agreement, four touches and safety.",WHITE,GREEN,18,14)
    add_card(slide,8.9,1.6,3.75,2.05,"3 · Semantic judge","Gemini evaluates support and relevance but cannot bypass a hard failure.",WHITE,AMBER,18,14)
    add_bullets(slide,["99% valid structured output", "95% citation coverage", "0 unsupported claims on golden cases", "100% outreach compliance and stop-rule tests"],.9,4.2,5.4,2.2,15,NAVY,VIOLET,.52)
    add_card(slide,6.65,4.2,5.7,1.65,"Two execution paths","Production calls providers and persists state. Evaluation uses frozen fixtures, avoiding Firecrawl/Apollo spend and CRM pollution.",VIOLET_PALE,VIOLET,17,13.5)
    footer(slide, 10)

    # 11 live data
    slide = base_slide(prs, "Q3 · Live data: the downstream control loop is proven", "Supabase rows · 10 Aug 2026")
    add_card(slide,.65,1.45,2.8,1.3,"3 synthetic contacts","Ranked 90 / 82 / 74",BLUE_PALE,BLUE,16,14)
    add_card(slide,3.7,1.45,2.8,1.3,"4 messages","1 sequence persisted",VIOLET_PALE,VIOLET,19,14)
    add_card(slide,6.75,1.45,2.8,1.3,"2 events","CTA click + Gmail reply",GREEN_PALE,GREEN,19,14)
    add_card(slide,9.8,1.45,2.8,1.3,"Stop applied","Sequence + peers paused",ORANGE_PALE,ORANGE,19,14)
    rows=[
        ("Jordan Lee", "Procurement Director", "90", "replied"),
        ("Taylor Morgan", "Head of Culinary & NPD", "82", "paused"),
        ("Alex Chen", "Operations Manager", "74", "paused"),
    ]
    add_table(slide,["CONTACT","ROLE","SCORE","STATE"],rows,.65,3.15,[2.55,4.0,1.4,2.0],.66,12)
    add_card(slide,10.85,3.15,1.85,2.64,"HONEST GAP","Newest research run remains running; older run is failed_partial. Live evidence and ingredient rows are not yet complete.",RED_PALE,RED,13,10.5)
    add_text(slide,"Live demo: https://knoxx-insight-quest.lovable.app",.75,6.25,7.0,.3,13,VIOLET,True)
    footer(slide, 11)

    # 12 screenshot/actionability
    slide = base_slide(prs, "What is accurate, actionable and still a hypothesis?", "Live proof vs. golden fixture")
    screenshot = DOCS / "account-intelligence.png"
    if screenshot.exists():
        slide.shapes.add_picture(str(screenshot), Inches(.65), Inches(1.45), width=Inches(7.1), height=Inches(4.65))
    add_card(slide,8.05,1.45,4.65,1.35,"ACTIONABLE NOW","Account/run identity, ranked contacts, safe sequence, click/reply events and organization stop.",GREEN_PALE,GREEN,16,12.5)
    add_card(slide,8.05,3.05,4.65,1.35,"USE WITH A LABEL","Scale, dishes, ingredient overlap and quantity ranges when linked to retained evidence.",AMBER_PALE,AMBER,16,12.5)
    add_card(slide,8.05,4.65,4.65,1.35,"SALES MUST VALIDATE","Service region, capacity, MOQ, margin, claims, recipe share and the real economic buyer.",RED_PALE,RED,16,12.5)
    add_text(slide,"The screenshot is a golden expected-output fixture, not live Firecrawl persistence proof.",.75,6.33,10.5,.32,12,SLATE,True)
    footer(slide, 12)

    # 13 scaling
    slide = base_slide(prs, "Q4 · What breaks at 1,000 domains?", "Bottlenecks and immediate controls")
    rows=[
        ("Crawl", "429s + long webhooks", "Async jobs, Retry-After, page bounds, cache"),
        ("LLM", "Cost + bad JSON/claims", "Bounded input, schema, repair once, escalate rarely"),
        ("Apollo", "Credits + plan limits", "Search first; enrich top 1–2; live quota checks"),
        ("n8n", "Bursts + long waits", "Queue mode, worker pools, provider backpressure"),
        ("Postgres", "Duplicates + contention", "Idempotency, short transactions, indexes"),
        ("Ops", "Stuck runs look green", "Checkpoints, traces, alerts, DLQ + replay"),
    ]
    add_table(slide,["BOTTLENECK","FAILURE MODE","CONTROL"],rows,.65,1.45,[2.0,3.55,6.5],.72,11.5)
    footer(slide, 13, "Firecrawl errors/crawl · Apollo rate limits/credits · n8n queue mode")

    # 14 enterprise target
    slide = base_slide(prs, "Enterprise target: accept → queue → checkpoint → notify", "Scale and cost")
    add_bullets(slide,[
        "Return 202 + run_id in under one second; never wait for crawling in the user request.",
        "Use durable jobs for retrieval, extraction, scoring, contacts and drafts—each with an idempotency key.",
        "Separate worker pools and rate budgets for Firecrawl, LLM and Apollo; retry transient failures, DLQ exhausted jobs.",
        "Cache evidence by content hash/freshness; reprocess only when evidence, catalogue or rule versions change.",
        "Measure cost per started, completed and qualified account—not an invented per-request promise.",
    ],.75,1.45,7.25,4.9,14.5,NAVY,ORANGE,.77)
    add_card(slide,8.35,1.45,4.0,2.0,"ILLUSTRATIVE CAPACITY","1,000 / 480 minutes ≈ 2.1 jobs/min. At 3 worker-minutes/job: ≈ 6.3 average concurrency; provision ~10–15 slots for variance.",NAVY,NAVY,16,14)
    add_card(slide,8.35,3.75,4.0,1.8,"STRONGEST COST CONTROL","If 30% reuse fresh evidence, plan 700 full runs + 300 revalidations—not 1,000 full crawls.",VIOLET_PALE,VIOLET,16,14)
    footer(slide, 14)

    # 15 close
    slide = base_slide(prs, "CPO decision: fund a controlled pilot", dark=True)
    add_bullets(slide,[
        "The product addresses account prioritization—not vanity automation.",
        "AI judgment is bounded by evidence, deterministic controls and human approval.",
        "The downstream click/reply/organization-stop loop has live proof.",
        "Production is blocked on live crawl persistence, sales policy, RLS/load and compliance checks.",
        "Next: a 25-account assisted pilot measuring speed, precision, trust and cost per accepted account.",
    ],.85,1.55,8.2,4.5,17,"FFFFFF",ORANGE,.77)
    add_card(slide,9.25,1.55,3.3,2.05,"PORTFOLIO HANDOFF","Architecture source\nGamma-ready deck\nLovable master prompt\nWorkflow + interview guide\nVerification report",NAVY_2,ORANGE,18,14)
    add_text(slide,"Questions",9.25,4.55,3.3,.55,28,WHITE,True)
    footer(slide, 15, dark=True)

    prs.save(OUT)
    print(OUT)
    print(ARCH_PNG)


if __name__ == "__main__":
    build_deck()
